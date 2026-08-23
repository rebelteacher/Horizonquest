"""Generate a HorizonQuest 'Import theme' .pptx (theme colors + fonts + dark master background)."""
import re, zipfile, shutil, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

MID = RGBColor(0x04, 0x12, 0x1F); TEAL = RGBColor(0x22, 0xD3, 0xEE)
ORANGE = RGBColor(0xFB, 0x92, 0x3C); WHITE = RGBColor(0xF7, 0xFA, 0xFC); MUTED = RGBColor(0x94, 0xA3, 0xB8)
HEAD, BODY = "Cormorant Garamond", "Outfit"
TMP = "/tmp/hq_theme_base.pptx"
OUT = "/app/frontend/public/decks/horizonquest_theme.pptx"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)


def bg(s, c=MID):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c


def txt(s, l, t, w, h, text, size, color, font, bold=False):
    tf = s.shapes.add_textbox(l, t, w, h).text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(size); r.font.name = font; r.font.bold = bold; r.font.color.rgb = color
    return tf


# Example slide 1 — title
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.6), Inches(0.18), Inches(1.8))
sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE; sh.line.fill.background(); sh.shadow.inherit = False
txt(s, Inches(1.2), Inches(2.3), Inches(11), Inches(1), "HorizonQuest Theme", 46, WHITE, HEAD, True)
txt(s, Inches(1.2), Inches(3.5), Inches(11), Inches(0.7), "Midnight blue · electric teal · sunrise orange", 20, TEAL, BODY)

# Example slide 2 — title + body
s2 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s2)
sh2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.6), Inches(0.16), Inches(0.7))
sh2.fill.solid(); sh2.fill.fore_color.rgb = TEAL; sh2.line.fill.background(); sh2.shadow.inherit = False
txt(s2, Inches(1.0), Inches(0.5), Inches(11.5), Inches(0.9), "Slide Title Style", 34, WHITE, HEAD, True)
tf = txt(s2, Inches(1.0), Inches(1.7), Inches(11.4), Inches(4.5), "This is body text in the theme font.", 20, WHITE, BODY)
for line, col in [("Accent bullet — teal marker", TEAL), ("Another key point", WHITE), ("Supporting detail", MUTED)]:
    p = tf.add_paragraph(); r = p.add_run(); r.text = "•  " + line
    r.font.size = Pt(18); r.font.name = BODY; r.font.color.rgb = col

prs.save(TMP)

# --- post-process theme1.xml + slideMaster background ---
CLR = ('<a:clrScheme name="HorizonQuest">'
       '<a:dk1><a:srgbClr val="04121F"/></a:dk1><a:lt1><a:srgbClr val="F7FAFC"/></a:lt1>'
       '<a:dk2><a:srgbClr val="0B1B30"/></a:dk2><a:lt2><a:srgbClr val="94A3B8"/></a:lt2>'
       '<a:accent1><a:srgbClr val="22D3EE"/></a:accent1><a:accent2><a:srgbClr val="FB923C"/></a:accent2>'
       '<a:accent3><a:srgbClr val="818CF8"/></a:accent3><a:accent4><a:srgbClr val="34D399"/></a:accent4>'
       '<a:accent5><a:srgbClr val="E11D48"/></a:accent5><a:accent6><a:srgbClr val="F59E0B"/></a:accent6>'
       '<a:hlink><a:srgbClr val="22D3EE"/></a:hlink><a:folHlink><a:srgbClr val="818CF8"/></a:folHlink>'
       '</a:clrScheme>')
FONTS = ('<a:fontScheme name="HorizonQuest">'
         '<a:majorFont><a:latin typeface="Cormorant Garamond"/><a:ea typeface=""/><a:cs typeface=""/></a:majorFont>'
         '<a:minorFont><a:latin typeface="Outfit"/><a:ea typeface=""/><a:cs typeface=""/></a:minorFont>'
         '</a:fontScheme>')
MASTER_BG = ('<p:bg><p:bgPr><a:solidFill><a:srgbClr val="04121F"/></a:solidFill><a:effectLst/></p:bgPr></p:bg>')

names = None
with zipfile.ZipFile(TMP, "r") as z:
    data = {n: z.read(n) for n in z.namelist()}

for name in list(data):
    if re.match(r"ppt/theme/theme\d+\.xml$", name):
        xml = data[name].decode("utf-8")
        xml = re.sub(r"<a:clrScheme.*?</a:clrScheme>", CLR, xml, count=1, flags=re.S)
        xml = re.sub(r"<a:fontScheme.*?</a:fontScheme>", FONTS, xml, count=1, flags=re.S)
        data[name] = xml.encode("utf-8")
    if re.match(r"ppt/slideMasters/slideMaster\d+\.xml$", name):
        xml = data[name].decode("utf-8")
        if "<p:bg>" in xml:
            xml = re.sub(r"<p:bg>.*?</p:bg>", MASTER_BG, xml, count=1, flags=re.S)
        else:
            xml = xml.replace("<p:cSld>", "<p:cSld>" + MASTER_BG, 1)
        data[name] = xml.encode("utf-8")

with zipfile.ZipFile(OUT, "w", zipfile.ZIP_DEFLATED) as z:
    for name, content in data.items():
        z.writestr(name, content)

os.remove(TMP)
print("Saved theme to", OUT, os.path.getsize(OUT), "bytes")
