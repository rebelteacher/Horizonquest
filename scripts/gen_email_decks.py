"""Generate themed HorizonQuest Email starter decks (.pptx) for import into Google Slides."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

MID = RGBColor(0x04, 0x12, 0x1F)
PANEL = RGBColor(0x0B, 0x1B, 0x30)
TEAL = RGBColor(0x22, 0xD3, 0xEE)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
WHITE = RGBColor(0xF7, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
HEAD = "Cormorant Garamond"
BODY = "Outfit"
W, H = Inches(13.333), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def bg(slide, color=MID):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    return tf


def run(p, text, size, color, font=BODY, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.name = font; r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return r


def title_slide(prs, kicker, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    rect(s, Inches(0.8), Inches(2.5), Inches(0.18), Inches(2.0), ORANGE)
    tf = tb(s, Inches(1.2), Inches(2.2), Inches(11), Inches(3))
    p = tf.paragraphs[0]; run(p, kicker, 16, TEAL, BODY, bold=True)
    p2 = tf.add_paragraph(); run(p2, title, 46, WHITE, HEAD, bold=True); p2.space_before = Pt(6)
    p3 = tf.add_paragraph(); run(p3, subtitle, 20, MUTED, BODY); p3.space_before = Pt(10)
    f = tb(s, Inches(1.2), Inches(6.7), Inches(11), Inches(0.5))
    run(f.paragraphs[0], "HorizonQuest · CTE Skill Studio", 12, TEAL, BODY, bold=True)


def content_slide(prs, title, bullets):
    """bullets: list of (text, level, kind) where kind in body/head/accent/muted."""
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    rect(s, Inches(0.7), Inches(0.6), Inches(0.16), Inches(0.7), TEAL)
    th = tb(s, Inches(1.0), Inches(0.5), Inches(11.6), Inches(0.9))
    run(th.paragraphs[0], title, 34, WHITE, HEAD, bold=True)
    body = tb(s, Inches(1.0), Inches(1.7), Inches(11.4), Inches(5.4))
    first = True
    for item in bullets:
        text, level, kind = (item + ("body",))[:3] if len(item) == 2 else item
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(6)
        if kind == "head":
            run(p, text, 22, TEAL, BODY, bold=True); p.space_before = Pt(8)
        elif kind == "accent":
            run(p, "▸ ", 18, ORANGE, BODY, bold=True); run(p, text, 18, WHITE, BODY)
        elif kind == "muted":
            run(p, text, 15, MUTED, BODY, italic=True)
        else:
            bullet = "•  " if level == 0 else "–  "
            run(p, bullet, 18, TEAL, BODY, bold=True)
            run(p, text, 18, WHITE if level == 0 else MUTED, BODY)
    return s


# ---------------- BLOCK 1 (from teacher notes) ----------------
d1 = new_deck()
title_slide(d1, "EMAIL · BLOCK 1", "Finding, Reading & Understanding Emails Like a Pro",
            "Locate senders, read the address lines, open emails, and search like a pro.")
content_slide(d1, "Learning Goals", [
    ("By the end of this lesson you will be able to:", 0, "muted"),
    ("Quickly find the right sender", 0, "accent"),
    ("Understand the difference between To, Cc, and Bcc", 0, "accent"),
    ("Locate and open any email", 0, "accent"),
    ("Identify all the important parts of an email", 0, "accent"),
    ("Use the search box to find emails in seconds", 0, "accent"),
    ("These skills help you stay organized at school, in clubs, and later at work!", 0, "muted"),
])
content_slide(d1, "1. Locate the Right Sender", [
    ("Every email has a From line that shows who sent it.", 0, "body"),
    ("How to find the right person", 0, "head"),
    ("Look at the From name or address at the top of the message", 0, "body"),
    ("In your inbox, the sender's name is usually bold on the left", 0, "body"),
    ("Seeing the same name a lot? Check the subject and date", 0, "body"),
    ("Pro tip: people use nicknames or school accounts — double-check the full address if it looks unfamiliar.", 0, "muted"),
])
content_slide(d1, "2. To, Cc & Bcc — What's the Difference?", [
    ("These three boxes control who receives the email and who can see each other.", 0, "muted"),
    ("To = Primary recipients — everyone can see them — the main people the email is for", 0, "body"),
    ("Cc = Carbon Copy — everyone can see them — people who need to know but don't have to reply", 0, "body"),
    ("Bcc = Blind Carbon Copy — only the sender & that person see it — keep an address private", 0, "body"),
    ("Easy way to remember", 0, "head"),
    ("To = the main people   ·   Cc = \"just so you know\" (visible)   ·   Bcc = secret copy (hidden)", 0, "accent"),
])
content_slide(d1, "To / Cc / Bcc — Example", [
    ("You email your teacher about a project:", 0, "body"),
    ("To: your teacher  (the main person)", 0, "accent"),
    ("Cc: your project partner  (needs to see it)", 0, "accent"),
    ("Bcc: your parent  (can see it without everyone knowing)", 0, "accent"),
])
content_slide(d1, "3. Locate and Open an Email", [
    ("Open your email app or website (Gmail, Outlook, school email)", 0, "body"),
    ("Go to your Inbox", 0, "body"),
    ("Scroll or look for the email you want", 0, "body"),
    ("Click or tap the email once — it opens and shows the full message", 0, "body"),
    ("Quick tip: unread emails are usually bold; once opened they become regular text.", 0, "muted"),
])
content_slide(d1, "4. Parts of an Email", [
    ("From → who sent it", 0, "body"),
    ("To / Cc / Bcc → who received it", 0, "body"),
    ("Date & Time → when it was sent", 0, "body"),
    ("Subject → short title telling you what it's about (very important!)", 0, "body"),
    ("Body → the actual message", 0, "body"),
    ("Attachments → files or pictures (look for a paperclip)", 0, "body"),
    ("Signature → the sender's name and contact info at the bottom", 0, "body"),
    ("Reply / Reply All / Forward → ways to respond or share", 0, "body"),
])
content_slide(d1, "5. Using the Search Box", [
    ("The search box is your superpower for finding emails fast (look for the magnifying glass).", 0, "muted"),
    ("from:ms.rivera  →  emails from Ms. Rivera", 0, "body"),
    ("subject:project  →  emails with \"project\" in the subject", 0, "body"),
    ("has:attachment  →  emails with files attached", 0, "body"),
    ("is:unread  →  only unread emails", 0, "body"),
    ("after:2026/08/01  →  emails after a date", 0, "body"),
    ("Combine words to get very specific results!", 0, "accent"),
])
content_slide(d1, "Quick Practice Challenges", [
    ("Find the Sender — write down a teacher/friend's full From address", 0, "body"),
    ("To/Cc/Bcc Check — open an email: who's in To? anyone in Cc/Bcc? why?", 0, "body"),
    ("Parts Hunt — label Subject, From, Date, Body, and Attachments", 0, "body"),
    ("Search Challenge — find an email from a person, one with an attachment, and an unread one", 0, "body"),
    ("Mental Map — teacher, partner, parent: which goes in To, Cc, or Bcc?", 0, "body"),
])
content_slide(d1, "Remember", [
    ("Always check the From and Subject first", 0, "accent"),
    ("To = main people · Cc = visible extras · Bcc = secret extras", 0, "accent"),
    ("The search box saves huge amounts of time", 0, "accent"),
    ("Knowing the parts of an email helps you read and reply confidently", 0, "accent"),
    ("You're now ready to handle email like a pro!", 0, "muted"),
])
d1.save("/app/frontend/public/decks/horizonquest_email_block1.pptx")

# ---------------- BLOCK 2: Replying, Forwarding & Composing ----------------
d2 = new_deck()
title_slide(d2, "EMAIL · BLOCK 2", "Replying, Forwarding & Composing",
            "Respond the right way, forward with a note, and write a complete email.")
content_slide(d2, "Learning Goals", [
    ("Reply to the sender vs. Reply All to everyone", 0, "accent"),
    ("Forward an email to someone new (with a short note)", 0, "accent"),
    ("Compose a brand-new email from scratch", 0, "accent"),
    ("Always follow email protocol: greeting, message, sign-off", 0, "accent"),
])
content_slide(d2, "Reply vs. Reply All", [
    ("Reply → sends your response to the original sender only", 0, "body"),
    ("Reply All → sends to the sender AND everyone else on the email", 0, "body"),
    ("The subject automatically starts with \"Re:\"", 0, "body"),
    ("Use Reply All only when everyone truly needs your answer", 0, "accent"),
    ("Reply-All-ing a whole class for a one-person answer spams everyone.", 0, "muted"),
])
content_slide(d2, "Forwarding an Email", [
    ("Forwarding = sending an email on to a new person", 0, "body"),
    ("The subject starts with \"Fwd:\"", 0, "body"),
    ("Put the new recipient in the To line", 0, "body"),
    ("Add a short note explaining why you're forwarding it", 0, "accent"),
])
content_slide(d2, "Composing a New Email", [
    ("Start a fresh message from scratch", 0, "body"),
    ("To: the main recipient's address", 0, "body"),
    ("Subject: short and clear about the topic (never leave it blank!)", 0, "body"),
    ("Body: your actual message", 0, "body"),
    ("Greeting + sign-off make it complete", 0, "accent"),
])
content_slide(d2, "Follow Email Protocol", [
    ("Every email you send should have three parts:", 0, "muted"),
    ("A greeting — e.g., \"Dear Ms. Lee,\" or \"Hi Coach,\"", 0, "accent"),
    ("A real message — a few clear sentences of your own", 0, "accent"),
    ("A sign-off — e.g., \"Thanks, Jordan\" or \"Sincerely, ...\"", 0, "accent"),
    ("Even a quick reply follows this — never send a blank email for credit.", 0, "muted"),
])
content_slide(d2, "Remember", [
    ("Re: = reply · Fwd: = forward", 0, "accent"),
    ("Reply = sender only · Reply All = everyone", 0, "accent"),
    ("Always greet, write a real message, and sign off", 0, "accent"),
])
d2.save("/app/frontend/public/decks/horizonquest_email_block2.pptx")

# ---------------- BLOCK 3: Recipients, Attachments, Formatting & Etiquette ----------------
d3 = new_deck()
title_slide(d3, "EMAIL · BLOCK 3", "Recipients, Attachments, Formatting & Etiquette",
            "Send like a professional — the right fields, files, formatting, and tone.")
content_slide(d3, "Learning Goals", [
    ("Choose the right recipient field: To, Cc, Bcc", 0, "accent"),
    ("Attach a file and mention it in your message", 0, "accent"),
    ("Use formatting (bold, bullets, signature) to be clear", 0, "accent"),
    ("Write with professional etiquette and tone", 0, "accent"),
])
content_slide(d3, "Choosing To / Cc / Bcc", [
    ("To → the main recipient(s)", 0, "body"),
    ("Cc → people who need a copy for awareness", 0, "body"),
    ("Bcc → copy someone privately (others can't see the address)", 0, "body"),
    ("Choosing the right field shows you understand who the message is for.", 0, "muted"),
])
content_slide(d3, "Attachments", [
    ("An attachment is a file sent along with the email (look for the paperclip)", 0, "body"),
    ("Common type for a document: PDF", 0, "body"),
    ("Always mention the attachment in your message", 0, "accent"),
    ("Forgot to attach it? The reader is left confused — double-check before sending.", 0, "muted"),
])
content_slide(d3, "Formatting Your Message", [
    ("Bold a key word to draw attention to important info", 0, "body"),
    ("Use a bulleted list when you have several items", 0, "body"),
    ("Add a signature with your name (and role/contact)", 0, "body"),
    ("Formatting should make the message clearer — not just decorate it.", 0, "muted"),
])
content_slide(d3, "Etiquette & Tone", [
    ("Be polite and professional", 0, "body"),
    ("WRITING IN ALL CAPS reads as shouting — avoid it", 0, "body"),
    ("Formal tone for a teacher, principal, or employer", 0, "body"),
    ("Informal tone is okay for a close friend or family member", 0, "body"),
    ("Proofread for tone, grammar, and spelling before you send.", 0, "accent"),
])
content_slide(d3, "Capstone: Manage Your Morning Inbox", [
    ("Put it all together in a real inbox:", 0, "muted"),
    ("Locate and open the right emails", 0, "body"),
    ("Reply and forward appropriately", 0, "body"),
    ("Compose a professional email with correct recipients", 0, "body"),
    ("Follow email protocol every time — greeting, message, sign-off.", 0, "accent"),
])
content_slide(d3, "Remember", [
    ("To = main · Cc = visible copy · Bcc = hidden copy", 0, "accent"),
    ("Mention every attachment you add", 0, "accent"),
    ("Professional tone + correct spelling = you look careful and capable", 0, "accent"),
])
d3.save("/app/frontend/public/decks/horizonquest_email_block3.pptx")

print("Saved 3 decks to /app/frontend/public/decks/")
